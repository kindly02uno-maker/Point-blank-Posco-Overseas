# -*- coding: utf-8 -*-
import sys, glob, os
from pptx import Presentation

def extract(path):
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"=== Slide {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in p.runs)
                    if t.strip():
                        out.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text for c in row.cells]
                    out.append(" | ".join(cells))
    return "\n".join(out)

files = glob.glob(sys.argv[1], recursive=True)
buf = []
for f in sorted(files):
    prs = Presentation(f)
    n = len(list(prs.slides))
    buf.append(f"\n\n########## FILE: {os.path.basename(f)} [slides: {n}] ##########")
    buf.append(extract(f))
with open(sys.argv[2], "w", encoding="utf-8") as fh:
    fh.write("\n".join(buf))
print("written to", sys.argv[2])
